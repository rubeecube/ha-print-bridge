"""Best-effort conversion of common printable documents to PDF.

The integration's print pipeline is PDF-centric.  This module keeps that
contract by converting common attachment formats into simple, readable PDFs
using Python-only dependencies.  Office fidelity is intentionally best-effort:
text, tables, sheets, slide text, and images are preserved where practical, but
exact pagination and styling are not guaranteed.
"""

from __future__ import annotations

import csv
import html
import io
import os
import re
from dataclasses import dataclass, field
from email.message import Message
from typing import Iterable

from PIL import Image, ImageSequence
from pypdf import PdfReader, PdfWriter

_PDF_MIME = "application/pdf"

_SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".docm",
    ".odt",
    ".rtf",
    ".txt",
    ".html",
    ".htm",
    ".md",
    ".xlsx",
    ".xlsm",
    ".xls",
    ".ods",
    ".csv",
    ".tsv",
    ".pptx",
    ".odp",
    ".jpg",
    ".jpeg",
    ".png",
    ".tif",
    ".tiff",
    ".bmp",
    ".webp",
}

_UNSUPPORTED_EXTENSIONS = {
    ".doc": "Legacy binary .doc files require LibreOffice or another office engine.",
    ".ppt": "Legacy binary .ppt files require LibreOffice or another office engine.",
}

_MIME_TO_EXTENSION = {
    "application/pdf": ".pdf",
    "application/msword": ".doc",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.ms-word.document.macroenabled.12": ".docm",
    "application/vnd.oasis.opendocument.text": ".odt",
    "application/rtf": ".rtf",
    "text/rtf": ".rtf",
    "text/plain": ".txt",
    "text/html": ".html",
    "text/markdown": ".md",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "application/vnd.ms-excel.sheet.macroenabled.12": ".xlsm",
    "application/vnd.ms-excel": ".xls",
    "application/vnd.oasis.opendocument.spreadsheet": ".ods",
    "text/csv": ".csv",
    "text/tab-separated-values": ".tsv",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.oasis.opendocument.presentation": ".odp",
    "image/jpeg": ".jpg",
    "image/png": ".png",
    "image/tiff": ".tiff",
    "image/bmp": ".bmp",
    "image/webp": ".webp",
}

_DOCUMENT_EXTENSIONS = {".docx", ".docm", ".odt", ".rtf", ".txt", ".html", ".htm", ".md"}
_SPREADSHEET_EXTENSIONS = {".xlsx", ".xlsm", ".xls", ".ods", ".csv", ".tsv"}
_PRESENTATION_EXTENSIONS = {".pptx", ".odp"}
_IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".tif", ".tiff", ".bmp", ".webp"}

SUPPORTED_PRINTABLE_EXTENSIONS = tuple(sorted(_SUPPORTED_EXTENSIONS))
UNSUPPORTED_PRINTABLE_EXTENSIONS = tuple(sorted(_UNSUPPORTED_EXTENSIONS))
SUPPORTED_PRINTABLE_MIME_TYPES = tuple(sorted(_MIME_TO_EXTENSION))


class UnsupportedDocumentError(ValueError):
    """Raised when a file type is known but cannot be converted safely."""


class DocumentConversionError(RuntimeError):
    """Raised when a supported file fails conversion."""


@dataclass(frozen=True)
class PrintableType:
    """Resolved printability metadata for one file."""

    extension: str
    source_format: str
    supported: bool
    reason: str | None = None


@dataclass(frozen=True)
class ConvertedDocument:
    """PDF bytes generated from one source document."""

    filename: str
    source_format: str
    pdf_data: bytes
    converted_format: str = _PDF_MIME


@dataclass
class AttachmentConversionSummary:
    """Summary of one combined message conversion."""

    converted: list[ConvertedDocument] = field(default_factory=list)
    skipped: list[str] = field(default_factory=list)

    @property
    def filenames(self) -> list[str]:
        return [document.filename for document in self.converted]


def normalise_content_type(content_type: str | None) -> str:
    """Return a bare lower-case MIME type."""
    return str(content_type or "").split(";", 1)[0].strip().lower()


def extension_for_document(filename: str | None, content_type: str | None = None) -> str:
    """Return the best extension hint from filename or MIME type."""
    ext = os.path.splitext(str(filename or ""))[1].lower()
    if ext:
        return ext
    mime = normalise_content_type(content_type)
    if mime.startswith("image/"):
        return _MIME_TO_EXTENSION.get(mime, ".png")
    return _MIME_TO_EXTENSION.get(mime, "")


def printable_type(filename: str | None, content_type: str | None = None) -> PrintableType | None:
    """Resolve whether an attachment is printable or a known unsupported format."""
    ext = extension_for_document(filename, content_type)
    mime = normalise_content_type(content_type)
    if ext in _SUPPORTED_EXTENSIONS:
        return PrintableType(
            extension=ext,
            source_format=_source_format(ext, mime),
            supported=True,
        )
    if ext in _UNSUPPORTED_EXTENSIONS:
        return PrintableType(
            extension=ext,
            source_format=_source_format(ext, mime),
            supported=False,
            reason=_UNSUPPORTED_EXTENSIONS[ext],
        )
    if mime.startswith("image/"):
        return PrintableType(
            extension=ext or ".image",
            source_format=mime,
            supported=True,
        )
    return None


def is_printable_attachment(
    filename: str | None,
    content_type: str | None = None,
    *,
    include_unsupported: bool = True,
) -> bool:
    """Return True when the file should be considered by print workflows."""
    resolved = printable_type(filename, content_type)
    if resolved is None:
        return False
    return resolved.supported or include_unsupported


def convert_document_to_pdf(
    data: bytes,
    filename: str,
    content_type: str | None = None,
) -> ConvertedDocument:
    """Convert supported source data to PDF bytes."""
    resolved = printable_type(filename, content_type)
    if resolved is None:
        raise UnsupportedDocumentError(f"Unsupported file type: {filename}")
    if not resolved.supported:
        raise UnsupportedDocumentError(resolved.reason or f"Unsupported file type: {filename}")

    ext = resolved.extension
    try:
        if ext == ".pdf":
            pdf_data = _validate_pdf(data)
        elif ext in _IMAGE_EXTENSIONS or resolved.source_format.startswith("image/"):
            pdf_data = _image_to_pdf(data, filename)
        elif ext in {".txt", ".html", ".htm", ".md", ".rtf"}:
            pdf_data = _textlike_to_pdf(data, filename, ext, content_type)
        elif ext in {".docx", ".docm"}:
            pdf_data = _docx_to_pdf(data, filename)
        elif ext == ".odt":
            pdf_data = _odt_to_pdf(data, filename)
        elif ext in {".csv", ".tsv"}:
            pdf_data = _csv_to_pdf(data, filename, ext, content_type)
        elif ext in {".xlsx", ".xlsm"}:
            pdf_data = _xlsx_to_pdf(data, filename)
        elif ext == ".xls":
            pdf_data = _xls_to_pdf(data, filename)
        elif ext == ".ods":
            pdf_data = _ods_to_pdf(data, filename)
        elif ext == ".pptx":
            pdf_data = _pptx_to_pdf(data, filename)
        elif ext == ".odp":
            pdf_data = _odp_to_pdf(data, filename)
        else:
            raise UnsupportedDocumentError(f"Unsupported file type: {filename}")
    except UnsupportedDocumentError:
        raise
    except Exception as exc:
        raise DocumentConversionError(f"Could not convert {filename}: {exc}") from exc

    return ConvertedDocument(
        filename=filename,
        source_format=resolved.source_format,
        pdf_data=pdf_data,
    )


def merge_pdf_documents(documents: Iterable[ConvertedDocument | bytes]) -> bytes:
    """Merge PDF byte streams into a single PDF."""
    writer = PdfWriter()
    count = 0
    for document in documents:
        pdf_data = document.pdf_data if isinstance(document, ConvertedDocument) else document
        reader = PdfReader(io.BytesIO(pdf_data))
        for page in reader.pages:
            writer.add_page(page)
            count += 1
    if count == 0:
        raise DocumentConversionError("No PDF pages were produced")
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def _source_format(ext: str, mime: str) -> str:
    if mime:
        return mime
    return ext.lstrip(".") or "unknown"


def _validate_pdf(data: bytes) -> bytes:
    PdfReader(io.BytesIO(data))
    return data


def _decode_text(data: bytes, content_type: str | None = None) -> str:
    charset = _charset_from_content_type(content_type)
    encodings = [charset, "utf-8-sig", "utf-16", "cp1252", "latin-1"]
    for encoding in [value for value in encodings if value]:
        try:
            return data.decode(encoding)
        except UnicodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _charset_from_content_type(content_type: str | None) -> str | None:
    if not content_type:
        return None
    msg = Message()
    msg["content-type"] = content_type
    return msg.get_content_charset()


def _textlike_to_pdf(
    data: bytes,
    filename: str,
    ext: str,
    content_type: str | None,
) -> bytes:
    text = _decode_text(data, content_type)
    if ext == ".rtf":
        from striprtf.striprtf import rtf_to_text

        text = rtf_to_text(text)
    elif ext == ".md":
        import markdown

        text = _html_to_text(markdown.markdown(text))
    elif ext in {".html", ".htm"}:
        text = _html_to_text(text)
    return _lines_to_pdf(filename, text.splitlines() or [""])


def _html_to_text(value: str) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(value, "html.parser")
    for element in soup(["script", "style"]):
        element.decompose()
    return html.unescape(soup.get_text("\n"))


def _docx_to_pdf(data: bytes, filename: str) -> bytes:
    from docx import Document

    document = Document(io.BytesIO(data))
    writer = _TextPdfWriter(title=filename)
    writer.heading(filename)
    for paragraph in document.paragraphs:
        writer.paragraph(paragraph.text)
    for table in document.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        writer.table(rows)
    return writer.finish()


def _odt_to_pdf(data: bytes, filename: str) -> bytes:
    from odf import table, text
    from odf.opendocument import load

    doc = load(io.BytesIO(data))
    writer = _TextPdfWriter(title=filename)
    writer.heading(filename)
    for paragraph in doc.getElementsByType(text.P):
        content = _odf_text(paragraph)
        if content:
            writer.paragraph(content)
    for odf_table in doc.getElementsByType(table.Table):
        writer.table(_odf_table_rows(odf_table))
    return writer.finish()


def _csv_to_pdf(data: bytes, filename: str, ext: str, content_type: str | None) -> bytes:
    text = _decode_text(data, content_type)
    delimiter = "\t" if ext == ".tsv" else ","
    rows = list(csv.reader(io.StringIO(text), delimiter=delimiter))
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    writer.heading(filename)
    writer.table(rows)
    return writer.finish()


def _xlsx_to_pdf(data: bytes, filename: str) -> bytes:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    writer.heading(filename)
    for sheet in workbook.worksheets:
        if getattr(sheet, "sheet_state", "visible") != "visible":
            continue
        writer.heading(sheet.title, level=2)
        writer.table(_normalise_table_rows(sheet.iter_rows(values_only=True)))
    return writer.finish()


def _xls_to_pdf(data: bytes, filename: str) -> bytes:
    import xlrd

    workbook = xlrd.open_workbook(file_contents=data)
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    writer.heading(filename)
    for sheet in workbook.sheets():
        writer.heading(sheet.name, level=2)
        rows = [[sheet.cell_value(row, col) for col in range(sheet.ncols)] for row in range(sheet.nrows)]
        writer.table(_normalise_table_rows(rows))
    return writer.finish()


def _ods_to_pdf(data: bytes, filename: str) -> bytes:
    from odf import table
    from odf.opendocument import load

    doc = load(io.BytesIO(data))
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    writer.heading(filename)
    for odf_table in doc.spreadsheet.getElementsByType(table.Table):
        writer.heading(str(odf_table.getAttribute("name") or "Sheet"), level=2)
        writer.table(_odf_table_rows(odf_table))
    return writer.finish()


def _pptx_to_pdf(data: bytes, filename: str) -> bytes:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    for index, slide in enumerate(presentation.slides, start=1):
        if index > 1:
            writer.new_page()
        writer.heading(f"{filename} - slide {index}")
        lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.extend(shape.text.splitlines())
        writer.lines(lines or ["(blank slide)"])
    return writer.finish()


def _odp_to_pdf(data: bytes, filename: str) -> bytes:
    from odf import draw, text
    from odf.opendocument import load

    doc = load(io.BytesIO(data))
    pages = doc.presentation.getElementsByType(draw.Page)
    writer = _TextPdfWriter(title=filename, landscape_page=True)
    for index, page in enumerate(pages, start=1):
        if index > 1:
            writer.new_page()
        writer.heading(f"{filename} - slide {index}")
        lines = [_odf_text(paragraph) for paragraph in page.getElementsByType(text.P)]
        writer.lines([line for line in lines if line] or ["(blank slide)"])
    return writer.finish()


def _image_to_pdf(data: bytes, filename: str) -> bytes:
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas

    image = Image.open(io.BytesIO(data))
    out = io.BytesIO()
    c = canvas.Canvas(out)
    margin = 36
    for index, frame in enumerate(ImageSequence.Iterator(image)):
        frame = _white_rgb_image(frame)
        page_size = landscape(A4) if frame.width >= frame.height else A4
        if index == 0:
            c.setPageSize(page_size)
        else:
            c.showPage()
            c.setPageSize(page_size)
        page_width, page_height = page_size
        max_width = page_width - (margin * 2)
        max_height = page_height - (margin * 2)
        scale = min(max_width / frame.width, max_height / frame.height)
        width = frame.width * scale
        height = frame.height * scale
        x = (page_width - width) / 2
        y = (page_height - height) / 2
        c.drawImage(ImageReader(frame), x, y, width=width, height=height)
    c.save()
    return out.getvalue()


def _white_rgb_image(image: Image.Image) -> Image.Image:
    image = image.convert("RGBA")
    background = Image.new("RGBA", image.size, "WHITE")
    background.alpha_composite(image)
    return background.convert("RGB")


def _lines_to_pdf(title: str, lines: list[str]) -> bytes:
    writer = _TextPdfWriter(title=title)
    writer.heading(title)
    writer.lines(lines)
    return writer.finish()


def _normalise_table_rows(rows: Iterable[Iterable[object]]) -> list[list[str]]:
    result: list[list[str]] = []
    for row in rows:
        values = ["" if value is None else str(value) for value in row]
        if any(value.strip() for value in values):
            result.append(values)
    return result or [[""]]


def _odf_text(node: object) -> str:
    parts: list[str] = []
    for child in getattr(node, "childNodes", []):
        if getattr(child, "nodeType", None) == child.TEXT_NODE:
            parts.append(str(getattr(child, "data", "")))
        else:
            parts.append(_odf_text(child))
    return "".join(parts).strip()


def _odf_table_rows(odf_table: object) -> list[list[str]]:
    from odf import table

    rows: list[list[str]] = []
    for row in odf_table.getElementsByType(table.TableRow):
        repeat_rows = int(row.getAttribute("numberrowsrepeated") or 1)
        values: list[str] = []
        for cell in row.getElementsByType(table.TableCell):
            repeat_cols = int(cell.getAttribute("numbercolumnsrepeated") or 1)
            value = _odf_text(cell)
            values.extend([value] * min(repeat_cols, 20))
        for _ in range(min(repeat_rows, 20)):
            rows.append(values)
    return _normalise_table_rows(rows)


class _TextPdfWriter:
    """Small reportlab wrapper for readable generated PDFs."""

    def __init__(self, *, title: str, landscape_page: bool = False) -> None:
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.pdfgen import canvas

        self._page_size = landscape(A4) if landscape_page else A4
        self._margin = 42
        self._out = io.BytesIO()
        self._canvas = canvas.Canvas(self._out, pagesize=self._page_size)
        self._font = _register_font()
        self._font_bold = self._font
        self._title = title
        self._y = self._page_size[1] - self._margin
        self._canvas.setTitle(title)

    @property
    def _width(self) -> float:
        return self._page_size[0]

    @property
    def _height(self) -> float:
        return self._page_size[1]

    @property
    def _content_width(self) -> float:
        return self._width - (self._margin * 2)

    def new_page(self) -> None:
        self._canvas.showPage()
        self._canvas.setPageSize(self._page_size)
        self._y = self._height - self._margin

    def heading(self, text: str, *, level: int = 1) -> None:
        size = 15 if level == 1 else 12
        self._ensure_space(size + 12)
        self._canvas.setFont(self._font_bold, size)
        for line in self._wrap(text, size):
            self._canvas.drawString(self._margin, self._y, line)
            self._y -= size + 4
        self._y -= 6

    def paragraph(self, text: str) -> None:
        if not text.strip():
            self._y -= 8
            return
        self.lines(text.splitlines())
        self._y -= 4

    def lines(self, lines: Iterable[str]) -> None:
        size = 10
        self._canvas.setFont(self._font, size)
        for source_line in lines:
            wrapped = self._wrap(str(source_line), size) or [""]
            for line in wrapped:
                self._ensure_space(size + 4)
                self._canvas.drawString(self._margin, self._y, line)
                self._y -= size + 4

    def table(self, rows: Iterable[Iterable[object]]) -> None:
        normalised = _normalise_table_rows(rows)
        if not normalised:
            return
        size = 8
        max_cols = min(max(len(row) for row in normalised), 8)
        if max_cols <= 0:
            return
        col_width = self._content_width / max_cols
        for row in normalised:
            cells = [str(value) for value in row[:max_cols]]
            wrapped_cells = [self._wrap(cell, size, max_width=col_width - 6) or [""] for cell in cells]
            row_height = (max(len(lines) for lines in wrapped_cells) * (size + 2)) + 8
            self._ensure_space(row_height)
            top = self._y
            for col_index, lines in enumerate(wrapped_cells):
                x = self._margin + (col_index * col_width)
                self._canvas.rect(x, top - row_height, col_width, row_height)
                text_y = top - size - 4
                self._canvas.setFont(self._font, size)
                for line in lines:
                    self._canvas.drawString(x + 3, text_y, line)
                    text_y -= size + 2
            self._y -= row_height
        self._y -= 8

    def finish(self) -> bytes:
        self._canvas.save()
        return self._out.getvalue()

    def _ensure_space(self, amount: float) -> None:
        if self._y - amount < self._margin:
            self.new_page()

    def _wrap(self, text: str, size: int, *, max_width: float | None = None) -> list[str]:
        from reportlab.pdfbase import pdfmetrics

        value = re.sub(r"\s+", " ", str(text or "")).strip()
        if not value:
            return [""]
        width = max_width or self._content_width
        words = value.split(" ")
        lines: list[str] = []
        current = ""
        for word in words:
            candidate = f"{current} {word}".strip()
            if pdfmetrics.stringWidth(candidate, self._font, size) <= width:
                current = candidate
                continue
            if current:
                lines.append(current)
            current = word
            while pdfmetrics.stringWidth(current, self._font, size) > width and len(current) > 1:
                split_at = max(1, int(len(current) * width / pdfmetrics.stringWidth(current, self._font, size)))
                lines.append(current[:split_at])
                current = current[split_at:]
        if current:
            lines.append(current)
        return lines


def _register_font() -> str:
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    for path in (
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/local/share/fonts/dejavu/DejaVuSans.ttf",
        "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/Library/Fonts/Arial.ttf",
    ):
        if not os.path.exists(path):
            continue
        try:
            pdfmetrics.registerFont(TTFont("PrintBridgeUnicode", path))
            return "PrintBridgeUnicode"
        except Exception:
            continue
    return "Helvetica"
