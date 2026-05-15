"""Tests for best-effort document attachment conversion."""

from __future__ import annotations

import io

import pytest
from pypdf import PdfReader, PdfWriter

from document_converter import (
    DocumentConversionError,
    UnsupportedDocumentError,
    convert_document_to_pdf,
    is_printable_attachment,
    merge_pdf_documents,
)


def _pdf_bytes(page_count: int = 1) -> bytes:
    writer = PdfWriter()
    for _ in range(page_count):
        writer.add_blank_page(width=595, height=842)
    out = io.BytesIO()
    writer.write(out)
    return out.getvalue()


def _page_count(data: bytes) -> int:
    return len(PdfReader(io.BytesIO(data)).pages)


@pytest.mark.parametrize(
    ("filename", "content_type", "expected"),
    [
        ("invoice.pdf", "application/octet-stream", True),
        ("letter.docx", "application/octet-stream", True),
        ("sheet.xlsx", "application/vnd.ms-excel", True),
        ("photo.jpg", "image/jpeg", True),
        ("legacy.doc", "application/msword", True),
        ("payload.json", "application/json", False),
    ],
)
def test_printable_detection_uses_mime_and_extension(
    filename: str, content_type: str, expected: bool
) -> None:
    assert is_printable_attachment(filename, content_type) is expected


def test_pdf_pass_through_and_merge() -> None:
    first = convert_document_to_pdf(_pdf_bytes(1), "one.pdf")
    second = convert_document_to_pdf(_pdf_bytes(2), "two.pdf")

    assert first.pdf_data.startswith(b"%PDF")
    assert first.source_format == "pdf"
    assert _page_count(merge_pdf_documents([first, second])) == 3


def test_text_html_markdown_and_rtf_convert_to_pdf() -> None:
    samples = [
        (b"hello\nworld", "notes.txt", None),
        (b"<h1>Hello</h1><p>world</p>", "page.html", "text/html"),
        (b"# Hello\n\nworld", "page.md", "text/markdown"),
        (br"{\rtf1\ansi Hello}", "page.rtf", "application/rtf"),
    ]
    for data, filename, content_type in samples:
        converted = convert_document_to_pdf(data, filename, content_type)
        assert _page_count(converted.pdf_data) >= 1
        assert converted.converted_format == "application/pdf"


def test_docx_csv_xlsx_pptx_and_image_convert_to_pdf() -> None:
    from docx import Document
    from openpyxl import Workbook
    from PIL import Image
    from pptx import Presentation

    doc = Document()
    doc.add_paragraph("Hello from docx")
    doc_buf = io.BytesIO()
    doc.save(doc_buf)

    workbook = Workbook()
    workbook.active["A1"] = "Hello"
    workbook.active["B1"] = "World"
    xlsx_buf = io.BytesIO()
    workbook.save(xlsx_buf)

    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Hello slide"
    pptx_buf = io.BytesIO()
    deck.save(pptx_buf)

    image_buf = io.BytesIO()
    Image.new("RGB", (80, 40), "white").save(image_buf, format="PNG")

    samples = [
        (doc_buf.getvalue(), "letter.docx"),
        (b"name,total\nA,12\n", "sheet.csv"),
        (xlsx_buf.getvalue(), "sheet.xlsx"),
        (pptx_buf.getvalue(), "deck.pptx"),
        (image_buf.getvalue(), "photo.png"),
    ]
    for data, filename in samples:
        converted = convert_document_to_pdf(data, filename)
        assert _page_count(converted.pdf_data) >= 1


def test_legacy_binary_office_files_are_rejected() -> None:
    with pytest.raises(UnsupportedDocumentError, match="Legacy binary .doc"):
        convert_document_to_pdf(b"not really word", "old.doc")

    with pytest.raises(UnsupportedDocumentError, match="Legacy binary .ppt"):
        convert_document_to_pdf(b"not really powerpoint", "old.ppt")


def test_corrupt_supported_file_reports_conversion_error() -> None:
    with pytest.raises(DocumentConversionError):
        convert_document_to_pdf(b"not a zip", "broken.docx")
